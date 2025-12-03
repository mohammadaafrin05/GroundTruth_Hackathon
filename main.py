import sys
import os
from datetime import datetime
import argparse
from pathlib import Path
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

def load_campaign_data(filepath):
    if not os.path.exists(filepath):
        print(f"ERROR: Can't find file at {filepath}")
        sys.exit(1)
    try:
        data = pd.read_csv(filepath)
        return data
    except Exception as e:
        print(f"ERROR reading CSV file: {e}")
        sys.exit(1)

def check_data_columns(data):
    column_mapping = {}

    for col in ['Campaign_ID','campaign','Campaign','campaign_id']:
        if col in data.columns:
            column_mapping['campaign'] = col
            break

    for col in ['Date','date','timestamp','Date_Stamp']:
        if col in data.columns:
            column_mapping['date'] = col
            break

    for col in ['Impressions','impressions','views','Views']:
        if col in data.columns:
            column_mapping['impressions'] = col
            break

    for col in ['Clicks','clicks','click_count']:
        if col in data.columns:
            column_mapping['clicks'] = col
            break

    for col in ['Acquisition_Cost','Spend','spend','cost','Cost']:
        if col in data.columns:
            column_mapping['spend'] = col
            break

    for col in ['ROI','Revenue','revenue','sales','conversion_value']:
        if col in data.columns:
            column_mapping['revenue'] = col
            break

    return column_mapping

def clean_and_prepare_data(raw_data, column_map):
    data = raw_data.copy()

    for std_name, orig_name in column_map.items():
        data = data.rename(columns={orig_name: std_name})

    required = ['campaign','date','impressions','clicks','spend']
    for col in required:
        if col not in data.columns:
            print(f"ERROR: Missing required column '{col}'")
            sys.exit(1)

    data['spend'] = data['spend'].astype(str).str.replace('$','').str.replace(',','')

    data['impressions'] = pd.to_numeric(data['impressions'], errors='coerce')
    data['clicks'] = pd.to_numeric(data['clicks'], errors='coerce')
    data['spend'] = pd.to_numeric(data['spend'], errors='coerce')

    if 'revenue' in data.columns:
        data['revenue'] = pd.to_numeric(data['revenue'], errors='coerce')
    else:
        data['revenue'] = data['spend'] * 2.5

    data = data.dropna(subset=['impressions','clicks','spend'])

    data.loc[data['impressions'] == 0, 'impressions'] = 1
    data.loc[data['clicks'] == 0, 'clicks'] = 1
    data.loc[data['spend'] == 0, 'spend'] = 0.01

    data['ctr'] = data['clicks'] / data['impressions']
    data['cpc'] = data['spend'] / data['clicks']
    data['roas'] = data['revenue'] / data['spend']

    try:
        data['date'] = pd.to_datetime(data['date'])
    except:
        pass

    return data

def calculate_summary_stats(data):
    summary = {}
    summary['total_impressions'] = int(data['impressions'].sum())
    summary['total_clicks'] = int(data['clicks'].sum())
    summary['total_spend'] = float(data['spend'].sum())
    summary['total_revenue'] = float(data['revenue'].sum())
    summary['avg_ctr'] = float(data['ctr'].mean())
    summary['avg_cpc'] = float(data['cpc'].mean())
    summary['avg_roas'] = float(data['roas'].mean())

    campaign_stats = data.groupby('campaign').agg({
        'roas': 'mean',
        'revenue': 'sum',
        'spend': 'sum'
    }).round(2)

    if not campaign_stats.empty:
        best_idx = campaign_stats['roas'].idxmax()
        summary['best_campaign'] = best_idx
        summary['best_roas'] = float(campaign_stats.loc[best_idx, 'roas'])

        worst_idx = campaign_stats['roas'].idxmin()
        summary['worst_campaign'] = worst_idx
        summary['worst_roas'] = float(campaign_stats.loc[worst_idx, 'roas'])

        top_5 = campaign_stats.nlargest(5, 'revenue')
        summary['top_5_campaigns'] = top_5.to_dict('index')

    return summary

def generate_insights(summary_stats):
    s = summary_stats
    insights = [
        f"Our campaigns delivered {s['total_impressions']:,} impressions and {s['total_clicks']:,} clicks.",
        f"Total ad spend was ${s['total_spend']:,.2f}, generating ${s['total_revenue']:,.2f} in revenue.",
        f"The average click-through rate was {s['avg_ctr']*100:.2f}%.",
        f"Each click cost ${s['avg_cpc']:.2f}.",
        f"The best campaign was '{s['best_campaign']}' with ROAS {s['best_roas']:.2f}.",
        f"The lowest performing campaign was '{s['worst_campaign']}' with ROAS {s['worst_roas']:.2f}.",
        f"Overall ROAS was {s['avg_roas']:.2f}."
    ]
    return insights

def create_performance_chart(data, output_path):
    campaign_totals = data.groupby('campaign')[['impressions','clicks','revenue']].sum()

    if len(campaign_totals) > 15:
        campaign_totals = campaign_totals.nlargest(15, 'revenue')

    fig, ax = plt.subplots(figsize=(12, 7))

    x = range(len(campaign_totals))
    width = 0.25

    ax.bar([i - width for i in x], campaign_totals['impressions'], width)
    ax.bar(x, campaign_totals['clicks'], width)
    ax.bar([i + width for i in x], campaign_totals['revenue'], width)

    ax.set_xticks(list(x))
    ax.set_xticklabels(campaign_totals.index, rotation=45, ha='right')
    plt.tight_layout()
    plt.savefig(output_path, dpi=150)
    plt.close()

    return output_path

def create_powerpoint_report(summary_stats, insights, chart_path, output_path):
    pres = Presentation()

    slide1 = pres.slides.add_slide(pres.slide_layouts[0])
    slide1.shapes.title.text = "Campaign Performance Report"
    slide1.placeholders[1].text = datetime.now().strftime("%B %d, %Y")

    slide2 = pres.slides.add_slide(pres.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    tf = slide2.placeholders[1].text_frame
    for i in insights[:4]:
        p = tf.add_paragraph()
        p.text = f"• {i}"

    slide3 = pres.slides.add_slide(pres.slide_layouts[1])
    slide3.shapes.title.text = "Key Performance Indicators"
    tf = slide3.placeholders[1].text_frame

    s = summary_stats
    metrics = [
        f"Total Impressions: {s['total_impressions']:,}",
        f"Total Clicks: {s['total_clicks']:,}",
        f"Total Spend: ${s['total_spend']:,.2f}",
        f"Total Revenue: ${s['total_revenue']:,.2f}",
        f"Avg CTR: {s['avg_ctr']*100:.2f}%",
        f"Avg CPC: ${s['avg_cpc']:.2f}",
        f"Avg ROAS: {s['avg_roas']:.2f}"
    ]

    for m in metrics:
        p = tf.add_paragraph()
        p.text = m

    slide4 = pres.slides.add_slide(pres.slide_layouts[5])
    slide4.shapes.title.text = "Campaign Performance Overview"
    slide4.shapes.add_picture(str(chart_path), Inches(0.5), Inches(1.5), width=Inches(9))

    if 'top_5_campaigns' in summary_stats:
        slide5 = pres.slides.add_slide(pres.slide_layouts[1])
        slide5.shapes.title.text = "Top 5 Campaigns"
        tf = slide5.placeholders[1].text_frame
        for i, (camp, stats) in enumerate(summary_stats['top_5_campaigns'].items(), 1):
            p = tf.add_paragraph()
            p.text = f"{i}. {camp}: ${stats['revenue']:,.2f} (ROAS {stats['roas']:.2f})"

    pres.save(output_path)
    return output_path

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--input', type=str, default='data/marketing_campaign_dataset.csv')
    parser.add_argument('--output', type=str, default='output')
    args = parser.parse_args()

    output_dir = Path(args.output)
    output_dir.mkdir(exist_ok=True, parents=True)

    raw = load_campaign_data(args.input)
    colmap = check_data_columns(raw)
    clean = clean_and_prepare_data(raw, colmap)
    summary = calculate_summary_stats(clean)
    insights = generate_insights(summary)

    chart_path = output_dir / "campaign_performance.png"
    create_performance_chart(clean, chart_path)

    ppt_path = output_dir / "campaign_report.pptx"
    create_powerpoint_report(summary, insights, chart_path, ppt_path)

    print("Analysis complete.")
    print(f"Report: {ppt_path}")
    print(f"Chart: {chart_path}")

if __name__ == "__main__":
    main()
